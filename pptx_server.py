# pptx_server.py
import io
import os
import tempfile
import requests
import matplotlib.pyplot as plt

from flask import Flask, request, send_file, jsonify
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from PIL import Image, ImageDraw, ImageFilter

# ---------- CONFIG ----------
TEMPLATE_PPTX_PATH = "/mnt/data/ACCINZIA (1).pptx"
FONT_NAME = os.environ.get("FONT_NAME", "Calibri")
TITLE_FONT_SIZE = 36
BODY_FONT_SIZE = 14

THEME_MAP = {
    "tech": ("#1F4E79", "#6FA8DC"),
    "food": ("#2E7D32", "#A5D6A7"),
    "finance": ("#0B1A34", "#6B8EA3"),
    "education": ("#6A1B9A", "#C39BD3"),
    "health": ("#EF6C00", "#FFCC80"),
    "default": ("#3E3E3E", "#BDBDBD"),
}

# Matplotlib config for polished charts
plt.rcParams["font.size"] = 12
plt.rcParams["savefig.transparent"] = True
plt.rcParams["axes.facecolor"] = "none"

app = Flask(__name__)

# ---------- HELPERS ----------
def hex_to_rgb(hex_color):
    hex_color = hex_color.lstrip("#")
    return tuple(int(hex_color[i:i+2], 16) for i in (0, 2, 4))

def pick_theme(company_name, business_type):
    key = (company_name + " " + business_type).lower()
    if any(k in key for k in ["tech", "ai", "saas", "software"]):
        return THEME_MAP["tech"]
    if any(k in key for k in ["food", "organic", "agri", "restaurant"]):
        return THEME_MAP["food"]
    if any(k in key for k in ["finance", "bank", "invest"]):
        return THEME_MAP["finance"]
    if any(k in key for k in ["education", "edtech", "school"]):
        return THEME_MAP["education"]
    if any(k in key for k in ["health", "wellness", "clinic"]):
        return THEME_MAP["health"]
    return THEME_MAP["default"]

def make_gradient_image(color_from, color_to, size=(1280, 720), vertical=True, blur=True):
    img = Image.new("RGB", size, color_from)
    draw = ImageDraw.Draw(img)
    w, h = size
    steps = h if vertical else w
    for i in range(steps):
        ratio = i / float(steps - 1) if steps > 1 else 0
        r = int(color_from[0] + (color_to[0] - color_from[0]) * ratio)
        g = int(color_from[1] + (color_to[1] - color_from[1]) * ratio)
        b = int(color_from[2] + (color_to[2] - color_from[2]) * ratio)
        if vertical:
            draw.line([(0, i), (w, i)], fill=(r, g, b))
        else:
            draw.line([(i, 0), (i, h)], fill=(r, g, b))
    if blur:
        img = img.filter(ImageFilter.GaussianBlur(radius=6))
    bio = io.BytesIO()
    img.save(bio, format="PNG")
    bio.seek(0)
    return bio

def set_slide_background_image(slide, img_bytes, prs):
    # Use presentation dimensions (correct API)
    slide.shapes.add_picture(
        img_bytes,
        Inches(0),
        Inches(0),
        width=prs.slide_width,
        height=prs.slide_height
    )

def fetch_image_bytes(url):
    try:
        r = requests.get(url, timeout=8)
        r.raise_for_status()
        return io.BytesIO(r.content)
    except Exception:
        return None

def safe_text(value):
    if value is None:
        return ""
    if isinstance(value, list):
        return "\n".join([str(x) for x in value])
    return str(value)

# ---------- Matplotlib chart helpers ----------
def make_premium_pie_chart(categories, values, colors=None, dpi=160, figsize=(5,5)):
    if not categories or not values:
        categories = ["N/A"]
        values = [1]
    fig, ax = plt.subplots(figsize=figsize, dpi=dpi)
    palette = colors or plt.cm.tab20.colors
    # ensure enough colors
    slice_colors = [palette[i % len(palette)] for i in range(len(values))]
    wedges, texts, autotexts = ax.pie(
        values,
        labels=categories,
        autopct="%1.1f%%",
        startangle=140,
        wedgeprops={"linewidth": 1, "edgecolor": "white"},
        colors=slice_colors,
        pctdistance=0.75,
    )
    # donut
    centre = plt.Circle((0, 0), 0.40, fc="white")
    ax.add_artist(centre)
    ax.axis("equal")
    # style autotexts
    for t in autotexts:
        t.set_color("black")
        t.set_fontsize(10)
    buf = io.BytesIO()
    plt.savefig(buf, format="png", transparent=True, bbox_inches="tight", pad_inches=0.1)
    buf.seek(0)
    plt.close(fig)
    return buf

# ---------- PPTX ASSEMBLY ----------
@app.route("/generate", methods=["POST"])
def generate():
    payload = request.get_json(force=True)
    if not payload:
        return jsonify({"error": "no json payload provided"}), 400

    company = payload.get("company_name", "Company")
    business_type = payload.get("nature_of_business", "")
    theme_from_hex, theme_to_hex = pick_theme(company, business_type)
    color_from = hex_to_rgb(theme_from_hex)
    color_to = hex_to_rgb(theme_to_hex)

    prs = Presentation()

    # try to fetch a logo from payload.logo_url or from TEMPLATE_PPTX_PATH
    logo_img_bytes = None
    if payload.get("logo_url"):
        logo_img_bytes = fetch_image_bytes(payload["logo_url"])

    if not logo_img_bytes and os.path.exists(TEMPLATE_PPTX_PATH):
        try:
            tpl = Presentation(TEMPLATE_PPTX_PATH)
            for s in tpl.slides:
                for shp in s.shapes:
                    # picture shape
                    if hasattr(shp, "image") and shp.image is not None:
                        logo_img_bytes = io.BytesIO(shp.image.blob)
                        logo_img_bytes.seek(0)
                        break
                if logo_img_bytes:
                    break
        except Exception:
            logo_img_bytes = None

    # helper to add a styled slide
    def add_styled_slide(title_text, body_lines=None, slide_index=0, chart=None):
        # build a subtle variant gradient per slide index
        start = tuple(int(color_from[i] + (color_to[i] - color_from[i]) * (slide_index / 12)) for i in range(3))
        end = tuple(int(color_from[i] + (color_to[i] - color_from[i]) * ((slide_index + 1) / 12)) for i in range(3))
        grad_img = make_gradient_image(start, end, size=(1280, 720), vertical=True, blur=True)

        slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout
        set_slide_background_image(slide, grad_img, prs)

        # Title
        left, top, width, height = Inches(0.5), Inches(0.3), Inches(9), Inches(1)
        title_box = slide.shapes.add_textbox(left, top, width, height)
        title_tf = title_box.text_frame
        title_tf.margin_top = 0
        title_tf.margin_bottom = 0
        title_tf.text = title_text or ""
        p0 = title_tf.paragraphs[0]
        p0.font.size = Pt(TITLE_FONT_SIZE)
        p0.font.bold = True
        p0.font.name = FONT_NAME
        p0.font.color.rgb = RGBColor(255, 255, 255)

        # logo top-right if exists
        if logo_img_bytes:
            try:
                # we need to copy bytes because add_picture will consume the stream
                logo_img_bytes.seek(0)
                logo_data = io.BytesIO(logo_img_bytes.read())
                logo_data.seek(0)
                slide.shapes.add_picture(logo_data, prs.slide_width - Inches(2.2), Inches(0.2), width=Inches(2))
                logo_img_bytes.seek(0)
            except Exception:
                pass

        # Body
        if body_lines:
            left, top, width, height = Inches(0.5), Inches(1.5), Inches(9), Inches(4.5)
            tb = slide.shapes.add_textbox(left, top, width, height)
            tf = tb.text_frame
            tf.word_wrap = True
            for i, line in enumerate(body_lines):
                if i == 0:
                    tf.text = line
                    p = tf.paragraphs[0]
                else:
                    p = tf.add_paragraph()
                    p.text = line
                p.level = 0
                p.font.size = Pt(BODY_FONT_SIZE)
                p.font.name = FONT_NAME
                p.font.color.rgb = RGBColor(255, 255, 255)

        # Chart: if chart dict provided, create matplotlib image and insert
        if chart:
            try:
                ctype = chart.get("type", "pie")
                categories = chart.get("categories", [])
                values = chart.get("values", [])
                # choose palette derived from theme
                palette = [tuple(x/255 for x in hex_to_rgb(theme_from_hex)), tuple(x/255 for x in hex_to_rgb(theme_to_hex))]
                # fallback: use theme hex strings for matplotlib (strings accepted)
                mpl_colors = [theme_from_hex, theme_to_hex, "#8888FF", "#66CC99", "#FFAA44"]
                chart_img = None
                if ctype in ("pie", "donut"):
                    chart_img = make_premium_pie_chart(categories, values, colors=mpl_colors)
                else:
                    # fallback simple pie for unsupported types in this helper
                    chart_img = make_premium_pie_chart(categories, values, colors=mpl_colors)
                if chart_img:
                    # insert picture centered-ish
                    slide.shapes.add_picture(chart_img, Inches(1.5), Inches(1.6), width=Inches(6))
            except Exception as e:
                print("chart render error:", e)

        return slide

    # Build slides (12)
    hero_lines = [safe_text(payload.get("tagline") or payload.get("company_name")), safe_text(payload.get("nature_of_business"))]
    add_styled_slide(payload.get("company_name", "Company"), body_lines=hero_lines, slide_index=0)

    overview = [safe_text(payload.get("nature_of_business", "")), safe_text(payload.get("short_description", ""))]
    add_styled_slide("Nature of Business", body_lines=overview, slide_index=1)

    vm_lines = [f"Vision: {safe_text(payload.get('vision',''))}", f"Mission: {safe_text(payload.get('mission',''))}"]
    add_styled_slide("Vision & Mission", body_lines=vm_lines, slide_index=2)

    problems = [f"• {p}" for p in payload.get("consumer_problems", [])] or ["No specific problems provided."]
    add_styled_slide("Problems Faced by Consumers", body_lines=problems, slide_index=3)

    solutions = [f"• {s}" for s in payload.get("solutions_provided", [])] or ["No solutions provided."]
    add_styled_slide("Solutions Provided", body_lines=solutions, slide_index=4)

    prods = [f"• {p}" for p in payload.get("products_services", [])] or ["No products/services listed."]
    add_styled_slide("Products & Services", body_lines=prods, slide_index=5)

    # Market share: accept list of {segment,percent} or dict
    market = payload.get("market_share", [])
    if isinstance(market, dict):
        market_items = [{"segment": k, "percent": v} for k, v in market.items()]
    else:
        market_items = market or []
    categories = [m.get("segment", "") for m in market_items] if market_items else ["N/A"]
    values = [float(m.get("percent", 0)) for m in market_items] if market_items else [100]
    add_styled_slide("Market Share", body_lines=[f"Total market coverage: {sum(values)}%"], slide_index=6,
                     chart={"type": "pie", "categories": categories, "values": values})

    target_lines = [f"• {t}" for t in payload.get("target_market", [])] or [payload.get("target_market", "Not specified")]
    add_styled_slide("Target Market", body_lines=target_lines, slide_index=7)

    usp_lines = [f"• {u}" for u in (payload.get("usp") or [])] or [payload.get("usp", "Not specified")]
    add_styled_slide("Unique Selling Proposition (USP)", body_lines=usp_lines, slide_index=8)

    contact = payload.get("company_contact", payload.get("contact_details", "N/A"))
    add_styled_slide("Contact Details", body_lines=[safe_text(contact)], slide_index=9)

    directors = payload.get("directors", [])
    if not directors and isinstance(payload.get("director_name"), str):
        directors = [{
            "name": payload.get("director_name"),
            "phone": payload.get("director_phone"),
            "email": payload.get("director_email"),
            "education": payload.get("director_qualification")
        }]
    dir_lines = []
    for d in directors:
        line = f"Name: {d.get('name','')}\nPhone: {d.get('phone','')}\nEmail: {d.get('email','')}\nEducation: {d.get('education','')}"
        dir_lines.append(line)
    add_styled_slide("Directors", body_lines=dir_lines or ["No director data provided."], slide_index=10)

    fund = payload.get("fund_deployment", {})
    if not fund:
        fund = {
            "testing_manufacturing": float(payload.get("fund_deployment_testing", 0)),
            "man_power": float(payload.get("fund_deployment_manpower", 0)),
            "outsourced_services": float(payload.get("fund_deployment_outsourced", 0)),
            "ip_costs": float(payload.get("fund_deployment_ip_costs", 0)),
            "travel": float(payload.get("fund_deployment_travel", 0)),
            "consumables": float(payload.get("fund_deployment_consumables", 0)),
            "contingency": float(payload.get("fund_deployment_contingency", 0)),
            "others": float(payload.get("fund_deployment_others", 0)),
        }
    categories_fd = list(fund.keys()) if fund else ["N/A"]
    values_fd = [float(v) for v in fund.values()] if fund else [100]
    add_styled_slide("Fund Deployment", body_lines=[f"Total: {sum(values_fd)}%"], slide_index=11,
                     chart={"type": "pie", "categories": categories_fd, "values": values_fd})

    # Save and return
    tmp = tempfile.NamedTemporaryFile(delete=False, suffix=".pptx")
    prs.save(tmp.name)
    tmp.flush()
    tmp.seek(0)
    return send_file(tmp.name, as_attachment=True, download_name=f"{company}_pitchdeck.pptx")


if __name__ == "__main__":
    # Production: gunicorn is used; this is for local testing
    app.run(host="0.0.0.0", port=5001)
